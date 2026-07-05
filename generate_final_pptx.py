import os
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_presentation():
    # Define paths
    pdf_pages_dir = "pdf_pages"
    output_path = r"C:\Users\VANDAN PATEL\.gemini\antigravity\scratch\LifeMate-AI\LifeMate_AI_Presentation_Final.pptx"
    brain_copy_path = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\bd150957-e9e4-4964-a43c-294775bf11af\LifeMate_AI_Presentation_Final.pptx"
    
    img_flow = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\a0f43dc2-132d-4bcf-95ab-2c47c8bff311\process_flow_diagram_1783236878041.png"
    img_mock = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\a0f43dc2-132d-4bcf-95ab-2c47c8bff311\dashboard_mockup_1783236891524.png"
    img_arch = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\a0f43dc2-132d-4bcf-95ab-2c47c8bff311\architecture_diagram_1783236904705.png"

    print("Initializing PowerPoint Presentation...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank layout

    # Helper function to add background and cover placeholder
    def prepare_slide(prs, page_num, add_cover=True):
        slide = prs.slides.add_slide(blank_layout)
        bg_image_path = os.path.join(pdf_pages_dir, f"page_{page_num}.png")
        
        # Add page image as background
        slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        
        # Add white cover rectangle to mask "Type here..."
        if add_cover:
            cover = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.5)
            )
            cover.fill.solid()
            cover.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cover.line.color.rgb = RGBColor(255, 255, 255) # Border is white (invisible)
            
        return slide

    # Helper function to style header/body paragraphs
    def add_section(tf, head, desc, idx):
        p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_head.text = head
        p_head.font.name = "Segoe UI"
        p_head.font.size = Pt(14)
        p_head.font.bold = True
        p_head.font.color.rgb = RGBColor(66, 133, 244) # Google Blue
        p_head.space_before = Pt(8) if idx > 0 else Pt(0)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = RGBColor(60, 64, 67) # Muted Dark Grey
        p_desc.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 1: Participant Details (Title Slide)
    # ----------------------------------------------------
    print("Generating Slide 1...")
    s1 = prepare_slide(prs, 1, add_cover=False)
    
    # Participant Name field
    tb_name = s1.shapes.add_textbox(Inches(2.5), Inches(5.08), Inches(9.5), Inches(0.4))
    tf_name = tb_name.text_frame
    tf_name.word_wrap = True
    p_name = tf_name.paragraphs[0]
    p_name.text = "Vandan Patel, Shivaxi Dave, Dhruv Shah"
    p_name.font.name = "Segoe UI"
    p_name.font.size = Pt(14)
    p_name.font.bold = True
    p_name.font.color.rgb = RGBColor(31, 31, 31)

    # Problem Statement field
    tb_prob = s1.shapes.add_textbox(Inches(2.5), Inches(5.68), Inches(9.5), Inches(1.2))
    tf_prob = tb_prob.text_frame
    tf_prob.word_wrap = True
    p_prob = tf_prob.paragraphs[0]
    p_prob.text = "Fragmentation of daily planners, health, budgeting, eco-habits, and emergency prep checkers makes everyday decision-making fragmented and inefficient."
    p_prob.font.name = "Arial"
    p_prob.font.size = Pt(12)
    p_prob.font.color.rgb = RGBColor(60, 64, 67)

    # ----------------------------------------------------
    # SLIDE 2: Brief about the idea
    # ----------------------------------------------------
    print("Generating Slide 2...")
    s2 = prepare_slide(prs, 2, add_cover=True)
    tb_s2 = s2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
    tf_s2 = tb_s2.text_frame
    tf_s2.word_wrap = True
    
    brief_items = [
        ("● Consolidates Essential Daily Services", "Instead of switching between multiple single-purpose apps, users access a single, highly integrated glassmorphic hub."),
        ("● Gemini AI at the Core", "Powers personalized, high-context recommendations tailored to user inputs (schedule, finances, health indicators)."),
        ("● Dual-Benefit Focus", "Helps individuals make smarter daily decisions (productivity, finance, wellness) while promoting social good (sustainability, community bonding, emergency preparedness)."),
        ("● Modern Interface", "Crafted with a premium dark-mode layout and responsive, glassmorphic card patterns for an engaging, easy-to-use user experience.")
    ]
    for idx, (head, desc) in enumerate(brief_items):
        add_section(tf_s2, head, desc, idx)

    # ----------------------------------------------------
    # SLIDE 3: Solution Explanation & Approach
    # ----------------------------------------------------
    print("Generating Slide 3...")
    s3 = prepare_slide(prs, 3, add_cover=True)
    tb_s3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
    tf_s3 = tb_s3.text_frame
    tf_s3.word_wrap = True
    
    approach_items = [
        ("● Google Cloud / Gemini AI Integration", 
         "We built LifeMate AI with a Python Flask backend integrated with the modern google-genai SDK. Using the low-latency, highly capable gemini-2.5-flash model, the application coordinates inputs from different everyday domains (planning, finance, health, ecology) through context-aware prompt templates to deliver structured markdown advice."),
        ("● Real-World Problem & Practical Impact", 
         "LifeMate AI targets app fatigue and cognitive fragmentation. Users get immediate value: personalized daily schedules, 50/30/20 budget diagnostics, health/nutrition plans, and neighborhood green-habits blueprints. It also serves as a life-saving helper by providing instant emergency Go-Bag lists and safety checklists for Flood, Fire, Earthquake, and Heatwaves."),
        ("● Core Workflow & Data Transformation", 
         "User Inputs -> Modular Payload Assembly -> Context-Engineered Prompt Compilation -> Gemini API Inference -> Responsive UI Markdown Rendering. It turns raw daily tasks, income/expenses, and age/lifestyle indicators into structured, actionable everyday decisions.")
    ]
    for idx, (head, desc) in enumerate(approach_items):
        add_section(tf_s3, head, desc, idx)

    # ----------------------------------------------------
    # SLIDE 4: Opportunities
    # ----------------------------------------------------
    print("Generating Slide 4...")
    s4 = prepare_slide(prs, 4, add_cover=True)
    tb_s4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
    tf_s4 = tb_s4.text_frame
    tf_s4.word_wrap = True
    
    opp_items = [
        ("● How different is it from any of the other existing ideas?",
         "Unlike single-purpose apps, LifeMate AI merges personal productivity (daily schedules, wellness, finance) with community sustainability and emergency preparedness in one glassmorphic hub. It replaces fragmented app-switching with form-based and conversational AI inputs."),
        ("● USP of the proposed solution",
         "A unified, conversational life assistant with pre-engineered, context-aware prompt templates, supporting session-level API key setups for secure developer and reviewer evaluation out of the box.")
    ]
    for idx, (head, desc) in enumerate(opp_items):
        add_section(tf_s4, head, desc, idx)

    # ----------------------------------------------------
    # SLIDE 5: List of Features Offered
    # ----------------------------------------------------
    print("Generating Slide 5...")
    s5 = prepare_slide(prs, 5, add_cover=True)
    tb_s5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.3))
    tf_s5 = tb_s5.text_frame
    tf_s5.word_wrap = True
    
    features = [
        ("🏠 Home Dashboard", "Modern, glassmorphic dashboard visualizing active modules status and shortcuts."),
        ("💬 AI Smart Assistant", "Users type any question (study guides, habits, advice) and get real-time answers."),
        ("📅 Daily Planner", "Takes tasks lists and available hours to formulate optimized schedules and productivity hacks."),
        ("💰 Budget Assistant", "Analyzes income/expenses, performs 50/30/20 diagnostics, and provides saving plans."),
        ("❤️ Wellness Guide", "Designs exercise plans, diet routines, and sleep routines adapted to user age and goals."),
        ("🌍 Community & Sustainability", "Generates waste-reduction guides and neighbor collaboration clean-up organizers."),
        ("🚨 Emergency Readiness", "Constructs emergency Go-Bag packing lists and safety checklists for Flood, Fire, Earthquake.")
    ]
    for idx, (head, desc) in enumerate(features):
        p = tf_s5.paragraphs[0] if idx == 0 else tf_s5.add_paragraph()
        p.text = head + ": "
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(66, 133, 244)
        p.space_before = Pt(4) if idx > 0 else Pt(0)
        
        run = p.add_run()
        run.text = desc
        run.font.name = "Arial"
        run.font.size = Pt(11)
        run.font.color.rgb = RGBColor(60, 64, 67)
        run.font.bold = False

    # ----------------------------------------------------
    # SLIDE 6: Process flow diagram
    # ----------------------------------------------------
    print("Generating Slide 6...")
    s6 = prepare_slide(prs, 6, add_cover=True)
    s6.shapes.add_picture(img_flow, Inches(1.5), Inches(1.8), Inches(10.33), Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 7: Wireframes/Mock diagrams
    # ----------------------------------------------------
    print("Generating Slide 7...")
    s7 = prepare_slide(prs, 7, add_cover=True)
    s7.shapes.add_picture(img_mock, Inches(1.5), Inches(1.8), Inches(10.33), Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 8: Architecture diagram
    # ----------------------------------------------------
    print("Generating Slide 8...")
    s8 = prepare_slide(prs, 8, add_cover=True)
    s8.shapes.add_picture(img_arch, Inches(1.5), Inches(1.8), Inches(10.33), Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 9: Technologies Used
    # ----------------------------------------------------
    print("Generating Slide 9...")
    s9 = prepare_slide(prs, 9, add_cover=True)
    tb_s9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.1))
    tf_s9 = tb_s9.text_frame
    tf_s9.word_wrap = True
    
    techs = [
        ("● Google Gemini API (gemini-2.5-flash)", "Core AI reasoning engine. We chose 2.5-flash for its low-latency performance, cost-efficiency, and advanced instruction-following capabilities across multiple topics."),
        ("● google-genai Python SDK", "Leverages the latest, officially-supported modern Google GenAI Client libraries interfaces for prompt completions requests."),
        ("● Flask Framework (Python)", "Serves as the server backbone, managing session-based credentials and routing endpoints."),
        ("● HTML5, CSS3 (Vanilla), JavaScript (ES6+)", "Client interface crafted with a responsive dark-mode Glassmorphism system, loading instantly with zero external dependencies."),
        ("● python-pptx & python-dotenv", "Enables programmatic generation of hackathon presentations and secure loading of system-level environment credentials.")
    ]
    for idx, (head, desc) in enumerate(techs):
        add_section(tf_s9, head, desc, idx)

    # ----------------------------------------------------
    # SLIDE 10: Snapshots of the prototype
    # ----------------------------------------------------
    print("Generating Slide 10...")
    s10 = prepare_slide(prs, 10, add_cover=True)
    s10.shapes.add_picture(img_mock, Inches(1.5), Inches(1.8), Inches(10.33), Inches(5.0))

    # ----------------------------------------------------
    # SLIDE 11: Thank You
    # ----------------------------------------------------
    print("Generating Slide 11...")
    prepare_slide(prs, 11, add_cover=False)

    # Save presentation
    print("Saving presentation...")
    prs.save(output_path)
    
    # Copy to brain folder for archiving
    os.makedirs(os.path.dirname(brain_copy_path), exist_ok=True)
    shutil.copy2(output_path, brain_copy_path)
    
    print(f"Presentation successfully updated and saved at:\n{output_path}\n{brain_copy_path}")

    # Attempt to copy to the main file name if it is unlocked
    main_path = r"C:\Users\VANDAN PATEL\.gemini\antigravity\scratch\LifeMate-AI\LifeMate_AI_Presentation.pptx"
    try:
        shutil.copy2(output_path, main_path)
        print(f"Also successfully saved to main file: {main_path}")
    except PermissionError:
        print(f"[Warning] Could not save to main file: {main_path} (File is currently open/locked by PowerPoint)")

if __name__ == "__main__":
    generate_presentation()
