import os
import shutil
import copy
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def update_presentation():
    # Define paths
    template_path = r"C:\Users\VANDAN PATEL\.gemini\antigravity\scratch\bah_presentation.pptx"
    output_path = r"C:\Users\VANDAN PATEL\.gemini\antigravity\scratch\LifeMate-AI\LifeMate_AI_Presentation_Updated.pptx"
    brain_copy_path = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\bd150957-e9e4-4964-a43c-294775bf11af\LifeMate_AI_Presentation_Updated.pptx"
    
    img_flow = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\a0f43dc2-132d-4bcf-95ab-2c47c8bff311\process_flow_diagram_1783236878041.png"
    img_mock = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\a0f43dc2-132d-4bcf-95ab-2c47c8bff311\dashboard_mockup_1783236891524.png"
    img_arch = r"C:\Users\VANDAN PATEL\.gemini\antigravity\brain\a0f43dc2-132d-4bcf-95ab-2c47c8bff311\architecture_diagram_1783236904705.png"

    print("Loading presentation template...")
    prs = pptx.Presentation(template_path)
    
    # ----------------------------------------------------
    # Helper to replace pictures (preserving background)
    # ----------------------------------------------------
    def replace_slide_picture(slide, new_image_path):
        for shape in list(slide.shapes):
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                if "Google Shape" in shape.name:
                    continue
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                shape.element.getparent().remove(shape.element)
                slide.shapes.add_picture(new_image_path, left, top, width, height)
                break

    # ----------------------------------------------------
    # SLIDE 1: Title & Details
    # ----------------------------------------------------
    print("Updating Slide 1...")
    s1 = prs.slides[0]
    for shape in s1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Team Name" in text:
                shape.text_frame.text = "Participant Name: Vandan Patel, Shivaxi Dave, Dhruv Shah"
            elif "Problem Statement" in text:
                shape.text_frame.text = "Problem Statement: Fragmentation of daily planners, health, budgeting, eco-habits, and emergency prep checkers makes everyday decision-making fragmented and inefficient."
            elif "Team Leader Name" in text:
                shape.text_frame.text = "Project Title: LifeMate AI"

    # ----------------------------------------------------
    # SLIDE 2: Brief about the idea
    # ----------------------------------------------------
    print("Updating Slide 2 (Brief about the idea)...")
    s2 = prs.slides[1]
    
    # Update title
    for shape in list(s2.shapes):
        if shape.has_text_frame and "Team Members" in shape.text_frame.text:
            shape.text_frame.text = "Brief about the idea"
            for p in shape.text_frame.paragraphs:
                p.font.name = "Google Sans"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 31, 31)
        elif shape.has_table:
            # Remove the table shape
            s2.shapes._spTree.remove(shape.element)

    # Copy body text box template from Slide 4 (index 3) to Slide 2 (index 1)
    s4_temp = prs.slides[3]
    body_shape_temp = None
    for shape in s4_temp.shapes:
        if shape.has_text_frame and "Multi-Source" in shape.text_frame.text:
            body_shape_temp = shape
            break

    if body_shape_temp:
        new_body_el = copy.deepcopy(body_shape_temp.element)
        s2.shapes._spTree.append(new_body_el)
        body_shape_s2 = s2.shapes[-1]
        
        tf = body_shape_s2.text_frame
        tf.clear()
        
        brief_bullets = [
            ("● Consolidates Essential Daily Services", "Instead of switching between multiple single-purpose apps, users access a single, highly integrated glassmorphic hub."),
            ("● Gemini AI at the Core", "Powers personalized, high-context recommendations tailored to user inputs (schedule, finances, health indicators)."),
            ("● Dual-Benefit Focus", "Helps individuals make smarter daily decisions (productivity, finance, wellness) while promoting social good (sustainability, community bonding, emergency preparedness)."),
            ("● Modern Interface", "Crafted with a premium dark-mode layout and responsive, glassmorphic card patterns for an engaging, easy-to-use user experience.")
        ]
        
        for idx, (head, desc) in enumerate(brief_bullets):
            p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p_head.text = head
            p_head.font.name = "Google Sans"
            p_head.font.size = Pt(13)
            p_head.font.bold = True
            p_head.font.color.rgb = RGBColor(66, 133, 244)
            p_head.space_before = Pt(6) if idx > 0 else Pt(0)
            
            p_desc = tf.add_paragraph()
            p_desc.text = desc
            p_desc.font.name = "Arial"
            p_desc.font.size = Pt(11)
            p_desc.font.color.rgb = RGBColor(60, 64, 67)
            p_desc.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 3: Solution Explanation & Approach
    # ----------------------------------------------------
    print("Updating Slide 3 (Solution Explanation & Approach)...")
    s3 = prs.slides[2]
    for shape in s3.shapes:
        if shape.has_text_frame and "Opportunity should be able to explain" in shape.text_frame.text:
            tf = shape.text_frame
            tf.clear()
            
            p_title = tf.paragraphs[0]
            p_title.text = "Your solution should be able to explain the following:"
            p_title.font.name = "Google Sans"
            p_title.font.size = Pt(22)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(31, 31, 31)
            p_title.space_after = Pt(12)
            
            explain_items = [
                ("● Google Cloud / Gemini AI Integration", 
                 "We built LifeMate AI with a Python Flask backend integrated with the modern google-genai SDK. Using the low-latency, highly capable gemini-2.5-flash model, the application coordinates inputs from different everyday domains (planning, finance, health, ecology) through context-aware prompt templates to deliver structured markdown advice."),
                ("● Real-World Problem & Practical Impact", 
                 "LifeMate AI targets app fatigue and cognitive fragmentation. Users get immediate value: personalized daily schedules, 50/30/20 budget diagnostics, health/nutrition plans, and neighborhood green-habits blueprints. It also serves as a life-saving helper by providing instant emergency Go-Bag lists and safety checklists for Flood, Fire, Earthquake, and Heatwaves."),
                ("● Core Workflow & Data Transformation", 
                 "User Inputs -> Modular Payload Assembly -> Context-Engineered Prompt Compilation -> Gemini API Inference -> Responsive UI Markdown Rendering. It turns raw daily tasks, income/expenses, and age/lifestyle indicators into structured, actionable everyday decisions.")
            ]
            
            for head, desc in explain_items:
                p_head = tf.add_paragraph()
                p_head.text = head
                p_head.font.name = "Google Sans"
                p_head.font.size = Pt(13)
                p_head.font.bold = True
                p_head.font.color.rgb = RGBColor(66, 133, 244)
                p_head.space_before = Pt(6)
                
                p_desc = tf.add_paragraph()
                p_desc.text = desc
                p_desc.font.name = "Arial"
                p_desc.font.size = Pt(10.5)
                p_desc.font.color.rgb = RGBColor(60, 64, 67)
                p_desc.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 4: Opportunities
    # ----------------------------------------------------
    print("Updating Slide 4 (Opportunities)...")
    s4 = prs.slides[3]
    title_shape = None
    body_shape = None
    for shape in s4.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "List of features offered" in text:
                title_shape = shape
            elif "Multi-Source" in text:
                body_shape = shape
                
    if title_shape:
        title_shape.text_frame.text = "Opportunities"
        for p in title_shape.text_frame.paragraphs:
            p.font.name = "Google Sans"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 31, 31)
            
    if body_shape:
        tf = body_shape.text_frame
        tf.clear()
        
        opps = [
            ("● How different is it from any of the other existing ideas?",
             "Unlike single-purpose apps, LifeMate AI merges personal productivity (daily schedules, wellness, finance) with community sustainability and emergency preparedness in one glassmorphic hub. It replaces fragmented app-switching with form-based and conversational AI inputs."),
            ("● USP of the proposed solution",
             "A unified, conversational life assistant with pre-engineered, context-aware prompt templates, supporting session-level API key setups for secure developer and reviewer evaluation out of the box.")
        ]
        
        for idx, (head, desc) in enumerate(opps):
            p_head = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p_head.text = head
            p_head.font.name = "Google Sans"
            p_head.font.size = Pt(13)
            p_head.font.bold = True
            p_head.font.color.rgb = RGBColor(66, 133, 244)
            p_head.space_before = Pt(10) if idx > 0 else Pt(0)
            
            p_desc = tf.add_paragraph()
            p_desc.text = desc
            p_desc.font.name = "Arial"
            p_desc.font.size = Pt(11)
            p_desc.font.color.rgb = RGBColor(60, 64, 67)
            p_desc.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 5: List of Features Offered (NEW SLIDE INSERTED)
    # ----------------------------------------------------
    print("Inserting Slide 5 (Features list)...")
    # Extract background from Slide 2 (index 1)
    temp_bg_path = "temp_bg.png"
    bg_extracted = False
    for shape in prs.slides[1].shapes:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE and "Google Shape" in shape.name:
            with open(temp_bg_path, "wb") as f:
                f.write(shape.image.blob)
            bg_extracted = True
            break
            
    # Add slide
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Reorder to index 4 (so it is Slide 5)
    sldIdLst = prs.slides._sldIdLst
    slide_id_element = sldIdLst[-1]
    sldIdLst.remove(slide_id_element)
    sldIdLst.insert(4, slide_id_element)
    
    # Add background image
    if bg_extracted:
        bg_pic = new_slide.shapes.add_picture(temp_bg_path, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        bg_pic.name = "Google Shape;NewSlideBackground;p16_5"
        
    # Copy textboxes from Slide 4 (currently index 3)
    s4_source = prs.slides[3]
    s5_dest = prs.slides[4]
    title_box_copied = None
    body_box_copied = None
    
    for shape in s4_source.shapes:
        if shape.has_text_frame:
            new_el = copy.deepcopy(shape.element)
            s5_dest.shapes._spTree.append(new_el)
            copied_shape = s5_dest.shapes[-1]
            if "Opportunities" in shape.text_frame.text:
                title_box_copied = copied_shape
            else:
                body_box_copied = copied_shape
                
    # Update text on Slide 5
    if title_box_copied:
        title_box_copied.text_frame.text = "List of features offered by the solution"
        for p in title_box_copied.text_frame.paragraphs:
            p.font.name = "Google Sans"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 31, 31)
            
    if body_box_copied:
        tf = body_box_copied.text_frame
        tf.clear()
        
        features = [
            ("🏠 Home Dashboard", "Modern, glassmorphic dashboard visualizing active modules status and shortcuts."),
            ("💬 AI Smart Assistant", "Users type any question (study guides, habits, advice) and get real-time answers."),
            ("📅 Daily Planner", "Takes tasks lists and available hours to formulate optimized schedules and productivity hacks."),
            ("💰 Budget Assistant", "Analyzes income/expenses, performs 50/30/20 diagnostics, and provides saving plans."),
            ("❤️ Wellness Guide", "Designs exercise plans, diet routines, and sleep routines adapted to user age and goals."),
            ("🌍 Community & Sustainability", "Generates waste-reduction guides and neighbor collaboration clean-up organizers."),
            ("🚨 Emergency Readiness", "Constructs emergency Go-Bag packing lists and safety checklists for Flood, Fire, Earthquake.")
        ]
        
        for idx, (head, desc) in enumerate(features):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = head + ": "
            p.font.name = "Google Sans"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(66, 133, 244)
            p.space_before = Pt(4) if idx > 0 else Pt(0)
            
            run = p.add_run()
            run.text = desc
            run.font.name = "Arial"
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(60, 64, 67)
            run.font.bold = False

    # Clean up temp background image
    if os.path.exists(temp_bg_path):
        os.remove(temp_bg_path)

    # ----------------------------------------------------
    # SLIDE 6: Process flow diagram
    # ----------------------------------------------------
    print("Updating Slide 6 (Process Flow)...")
    s6 = prs.slides[5]
    for shape in s6.shapes:
        if shape.has_text_frame and "Process Flow Diagram" in shape.text_frame.text:
            shape.text_frame.text = "Process flow diagram or Use-case diagram"
            for p in shape.text_frame.paragraphs:
                p.font.name = "Google Sans"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 31, 31)
                
    replace_slide_picture(s6, img_flow)

    # ----------------------------------------------------
    # SLIDE 7: Wireframes/Mock diagrams
    # ----------------------------------------------------
    print("Updating Slide 7 (Mockup)...")
    s7 = prs.slides[6]
    for shape in s7.shapes:
        if shape.has_text_frame and "Mockup" in shape.text_frame.text:
            shape.text_frame.text = "Wireframes/Mock diagrams of the proposed solution"
            for p in shape.text_frame.paragraphs:
                p.font.name = "Google Sans"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 31, 31)
                
    replace_slide_picture(s7, img_mock)

    # ----------------------------------------------------
    # SLIDE 8: Architecture diagram
    # ----------------------------------------------------
    print("Updating Slide 8 (Architecture)...")
    s8 = prs.slides[7]
    for shape in s8.shapes:
        if shape.has_text_frame and "Architecture" in shape.text_frame.text:
            shape.text_frame.text = "Architecture diagram of the proposed solution:"
            for p in shape.text_frame.paragraphs:
                p.font.name = "Google Sans"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 31, 31)
                
    replace_slide_picture(s8, img_arch)

    # ----------------------------------------------------
    # SLIDE 9: Technologies
    # ----------------------------------------------------
    print("Updating Slide 9 (Technologies)...")
    s9 = prs.slides[8]
    for shape in s9.shapes:
        if shape.has_text_frame and "Technologies" in shape.text_frame.text:
            tf = shape.text_frame
            tf.clear()
            
            p_title = tf.paragraphs[0]
            p_title.text = "Technologies / Google / Nvidia Services used in the solution"
            p_title.font.name = "Google Sans"
            p_title.font.size = Pt(22)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(31, 31, 31)
            p_title.space_after = Pt(12)
            
            techs = [
                ("● Google Gemini API (gemini-2.5-flash)", "Core AI reasoning engine. We chose 2.5-flash for its low-latency performance, cost-efficiency, and advanced instruction-following capabilities across multiple topics."),
                ("● google-genai Python SDK", "Leverages the latest, officially-supported modern Google GenAI Client libraries interfaces for prompt completions requests."),
                ("● Flask Framework (Python)", "Serves as the server backbone, managing session-based credentials and routing endpoints."),
                ("● HTML5, CSS3 (Vanilla), JavaScript (ES6+)", "Client interface crafted with a responsive dark-mode Glassmorphism system, loading instantly with zero external dependencies."),
                ("● python-pptx & python-dotenv", "Enables programmatic generation of hackathon presentations and secure loading of system-level environment credentials.")
            ]
            
            for head, desc in techs:
                p_head = tf.add_paragraph()
                p_head.text = head
                p_head.font.name = "Google Sans"
                p_head.font.size = Pt(13)
                p_head.font.bold = True
                p_head.font.color.rgb = RGBColor(66, 133, 244)
                p_head.space_before = Pt(6)
                
                p_desc = tf.add_paragraph()
                p_desc.text = desc
                p_desc.font.name = "Arial"
                p_desc.font.size = Pt(11)
                p_desc.font.color.rgb = RGBColor(60, 64, 67)
                p_desc.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 10: Snapshots of the prototype
    # ----------------------------------------------------
    print("Updating Slide 10 (Snapshots of prototype)...")
    s10 = prs.slides[9]
    table_shape = None
    for shape in list(s10.shapes):
        if shape.has_text_frame and "Cost" in shape.text_frame.text:
            shape.text_frame.text = "Snapshots of the prototype"
            for p in shape.text_frame.paragraphs:
                p.font.name = "Google Sans"
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(31, 31, 31)
        elif shape.has_table:
            table_shape = shape
            
    if table_shape:
        s10.shapes._spTree.remove(table_shape.element)
        
    # Add prototype image
    s10.shapes.add_picture(img_mock, Inches(0.75), Inches(1.575), Inches(8.5), Inches(4.5))

    # ----------------------------------------------------
    # SLIDE 11: Thank You
    # ----------------------------------------------------
    print("Slide 11 (Thank you) preserved.")

    # Save presentation
    print("Saving presentation...")
    prs.save(output_path)
    
    # Copy to brain folder for archiving
    os.makedirs(os.path.dirname(brain_copy_path), exist_ok=True)
    shutil.copy2(output_path, brain_copy_path)
    
    print(f"Presentation successfully updated and saved at:\n{output_path}\n{brain_copy_path}")

if __name__ == "__main__":
    update_presentation()
