import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    COLOR_BG = RGBColor(12, 14, 26)       # Dark Navy
    COLOR_CYAN = RGBColor(0, 242, 254)     # Glowing Cyan
    COLOR_PURPLE = RGBColor(177, 159, 251) # Soft Purple
    COLOR_WHITE = RGBColor(243, 244, 246)  # Off-White
    COLOR_MUTED = RGBColor(156, 163, 175)  # Muted Grey
    COLOR_ORANGE = RGBColor(255, 140, 0)   # Accent Orange

    # Helper function to set solid background color
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    # Helper function to add a title
    def add_slide_title(slide, text, color=COLOR_CYAN):
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Outfit'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        return title_box

    # Blank layout is the 6th layout (index 6)
    blank_layout = prs.slide_layouts[6]

    # ==========================================
    # SLIDE 1: Title & Participant Details
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_1)

    # Glowing title box in center-left
    title_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "LifeMate AI"
    p.font.name = 'Outfit'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    p2 = tf.add_paragraph()
    p2.text = "Your AI companion for smarter living and stronger communities."
    p2.font.name = 'Inter'
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_PURPLE
    p2.space_before = Pt(12)

    # Participant Details Box
    details_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.0))
    tf_details = details_box.text_frame
    tf_details.word_wrap = True
    tf_details.margin_left = tf_details.margin_right = tf_details.margin_top = tf_details.margin_bottom = 0

    dp1 = tf_details.paragraphs[0]
    dp1.text = "PARTICIPANT NAME:"
    dp1.font.name = 'Outfit'
    dp1.font.size = Pt(14)
    dp1.font.bold = True
    dp1.font.color.rgb = COLOR_MUTED

    dp2 = tf_details.add_paragraph()
    dp2.text = "Vandan Patel"
    dp2.font.name = 'Inter'
    dp2.font.size = Pt(18)
    dp2.font.bold = True
    dp2.font.color.rgb = COLOR_WHITE
    dp2.space_after = Pt(16)

    dp3 = tf_details.add_paragraph()
    dp3.text = "PROBLEM STATEMENT:"
    dp3.font.name = 'Outfit'
    dp3.font.size = Pt(14)
    dp3.font.bold = True
    dp3.font.color.rgb = COLOR_MUTED
    dp3.space_before = Pt(8)

    dp4 = tf_details.add_paragraph()
    dp4.text = "People use different apps for planning, health, budgeting, eco-habits, and emergency guidance. This makes everyday decision-making fragmented and inefficient."
    dp4.font.name = 'Inter'
    dp4.font.size = Pt(16)
    dp4.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 2: Brief about the idea
    # ==========================================
    slide_2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_2)
    add_slide_title(slide_2, "Brief about the idea")

    content_box = slide_2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.text = "LifeMate AI: A Unified Companion for Smarter Living"
    p.font.name = 'Outfit'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PURPLE
    p.space_after = Pt(20)

    bullets = [
        "Consolidates essential daily services: Instead of switching between multiple single-purpose apps, users access a single, highly integrated hub.",
        "Gemini AI at the Core: Powers personalized, high-context recommendations tailored to user inputs (schedule, finances, health indicators).",
        "Dual-Benefit Focus: Helps individuals make smarter daily decisions (productivity, finance, wellness) while promoting social good (sustainability, community bonding, emergency preparedness).",
        "Modern Interface: Crafted with a premium dark-mode layout and responsive, glassmorphic card patterns for an engaging, easy-to-use user experience."
    ]

    for bullet in bullets:
        bp = tf.add_paragraph()
        bp.text = "• " + bullet.split(":")[0] + ":"
        bp.font.name = 'Inter'
        bp.font.size = Pt(18)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_CYAN
        bp.space_before = Pt(10)
        
        desc = bullet.split(":")[1]
        bp_desc = tf.add_paragraph()
        bp_desc.text = desc.strip()
        bp_desc.font.name = 'Inter'
        bp_desc.font.size = Pt(16)
        bp_desc.font.color.rgb = COLOR_WHITE
        bp_desc.space_after = Pt(12)
        bp_desc.level = 0
        bp_desc.space_before = Pt(2)

    # ==========================================
    # SLIDE 3: Solution Explanation & Approach
    # ==========================================
    slide_3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_3)
    add_slide_title(slide_3, "Your solution should be able to explain the following:")

    col_width = Inches(3.7)
    gap = Inches(0.3)
    left_start = Inches(0.8)
    top_pos = Inches(1.8)

    questions = [
        ("Approach & Google Tech", 
         "How did you approach the problem and translate it into a working solution using Google Cloud?",
         ["Leveraged Google Gemini API to coordinate multi-domain inputs.",
          "Used Flask as a lightweight backend coordinator for context prompts.",
          "Configured API Key setup inside user sessions, keeping credentials secure."]),
        ("Real-World Impact", 
         "What real-world problem does your solution address, and what practical impact does it create?",
         ["Combats cognitive fatigue caused by fragmented app layouts.",
          "Empowers local action (composting hubs, tool sharing).",
          "Distributes immediate safety checklists for Floods/Fires/Earthquakes."]),
        ("Core Workflow / Architecture", 
         "What is the core workflow behind your solution, and how does it transform data into insights?",
         ["Inputs -> Prompt engineering constructs tailored templates.",
          "Gemini processes inputs & returns markdown schedules/checklists.",
          "Frontend parses markdown dynamically for responsive UI display."])
    ]

    for idx, (title, sub, items) in enumerate(questions):
        box_left = left_start + idx * (col_width + gap)
        # Background panels
        shape = slide_3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, box_left, top_pos, col_width, Inches(5.0)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(18, 20, 36)
        shape.line.color.rgb = COLOR_PURPLE
        shape.line.width = Pt(1.5)

        # Text Frame
        tb = slide_3.shapes.add_textbox(box_left + Inches(0.2), top_pos + Inches(0.2), col_width - Inches(0.4), Inches(4.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Outfit'
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.space_after = Pt(10)

        p_sub = tf.add_paragraph()
        p_sub.text = sub
        p_sub.font.name = 'Inter'
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = COLOR_MUTED
        p_sub.font.italic = True
        p_sub.space_after = Pt(14)

        for item in items:
            ip = tf.add_paragraph()
            ip.text = "• " + item
            ip.font.name = 'Inter'
            ip.font.size = Pt(14)
            ip.font.color.rgb = COLOR_WHITE
            ip.space_before = Pt(8)

    # ==========================================
    # SLIDE 4: Opportunities & USP
    # ==========================================
    slide_4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_4)
    add_slide_title(slide_4, "Opportunities")

    tb_opp = slide_4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_opp = tb_opp.text_frame
    tf_opp.word_wrap = True
    tf_opp.margin_left = tf_opp.margin_right = tf_opp.margin_top = tf_opp.margin_bottom = 0

    op1 = tf_opp.paragraphs[0]
    op1.text = "How different is it from other existing ideas?"
    op1.font.name = 'Outfit'
    op1.font.size = Pt(22)
    op1.font.bold = True
    op1.font.color.rgb = COLOR_CYAN
    op1.space_after = Pt(14)

    opp_bullets = [
        "Multi-domain correlation: Integrates calendar scheduler directly with wellness habits and budget guides, allowing users to save time and energy.",
        "Proactive Community support: Rather than just tracking personal habits, it coordinates local sustainable plans and neighbors clean-ups.",
        "Interactive Emergency Readiness: Provides context-specific, localized guides (e.g. apartment vs. residential house safety) for severe hazards."
    ]

    for bullet in opp_bullets:
        p = tf_opp.add_paragraph()
        p.text = "• " + bullet.split(":")[0] + ":"
        p.font.name = 'Inter'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PURPLE
        p.space_before = Pt(8)
        
        p_desc = tf_opp.add_paragraph()
        p_desc.text = bullet.split(":")[1].strip()
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = COLOR_WHITE
        p_desc.space_after = Pt(8)

    # USP Column
    tb_usp = slide_4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_usp = tb_usp.text_frame
    tf_usp.word_wrap = True
    tf_usp.margin_left = tf_usp.margin_right = tf_usp.margin_top = tf_usp.margin_bottom = 0

    up1 = tf_usp.paragraphs[0]
    up1.text = "USP of the proposed solution"
    up1.font.name = 'Outfit'
    up1.font.size = Pt(22)
    up1.font.bold = True
    up1.font.color.rgb = COLOR_ORANGE
    up1.space_after = Pt(14)

    usp_bullets = [
        "All-in-One Life Dashboard: Unified console reducing app context switching and decision fragmentation.",
        "Context-Aware Prompt Engineering: Pre-engineered high-context prompt templates optimized for Gemini API inference.",
        "Robust Double SDK Compatibility: Backend supports both legacy Google GenerativeAI and modern Google GenAI Python SDK wrappers.",
        "Reviewer-Friendly Setup: Session-based API key input avoids credentials leaks while making the app immediately reviewable out of the box."
    ]

    for bullet in usp_bullets:
        p = tf_usp.add_paragraph()
        p.text = "✔ " + bullet.split(":")[0] + ":"
        p.font.name = 'Inter'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.space_before = Pt(8)
        
        p_desc = tf_usp.add_paragraph()
        p_desc.text = bullet.split(":")[1].strip()
        p_desc.font.name = 'Inter'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = COLOR_MUTED
        p_desc.space_after = Pt(8)

    # ==========================================
    # SLIDE 5: List of Features
    # ==========================================
    slide_5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_5)
    add_slide_title(slide_5, "List of features offered by the solution")

    tb_feat = slide_5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_feat = tb_feat.text_frame
    tf_feat.word_wrap = True
    tf_feat.margin_left = tf_feat.margin_right = tf_feat.margin_top = tf_feat.margin_bottom = 0

    features = [
        ("🏠 Home Dashboard", "Modern, glassmorphic dashboard visualizing active modules status and shortcuts to tools."),
        ("📅 Daily Planner", "Takes task list, available hours, and wellness goals to formulate an optimized hourly schedule and productivity hacks."),
        ("💰 Budget Assistant", "Computes income vs. expenses, breaks down finance using the 50/30/20 rule, and details specific savings action items."),
        ("❤️ Wellness Guide", "Builds exercise plans, diet ideas, and sleep routines adapted to the user's age, fitness goals, and dietary restrictions."),
        ("🌍 Community & Eco Hub", "Constructs green initiatives plans, composting setups, and community cleanup organizers based on user interest."),
        ("🚨 Emergency Center", "Generates immediate survival actions checklists and emergency Go-Bag packers for Flood, Fire, Earthquake, and Heatwaves."),
        ("💬 AI Smart Assistant", "General chat assistant answering random daily questions like study guides, exam preparation, and smart habits.")
    ]

    for idx, (title, desc) in enumerate(features):
        p = tf_feat.paragraphs[0] if idx == 0 else tf_feat.add_paragraph()
        p.text = title + ": "
        p.font.name = 'Outfit'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.space_before = Pt(8)
        
        run = p.add_run()
        run.text = desc
        run.font.name = 'Inter'
        run.font.size = Pt(15)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = False

    # ==========================================
    # SLIDE 6: Process flow diagram or Use-case
    # ==========================================
    slide_6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_6)
    add_slide_title(slide_6, "Process flow diagram or Use-case diagram")

    tb_flow = slide_6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_flow = tb_flow.text_frame
    tf_flow.word_wrap = True
    tf_flow.margin_left = tf_flow.margin_right = tf_flow.margin_top = tf_flow.margin_bottom = 0

    p = tf_flow.paragraphs[0]
    p.text = "Application Core Logic & Data Flow"
    p.font.name = 'Outfit'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(20)

    # Drawing a text-based process map
    flow_steps = [
        ("1. User Inputs & Configs", "User interacts with modular forms or chat interface. Enters API Key securely via settings panel (saved in Flask session)."),
        ("2. Backend Routing", "Flask router `/api/generate` intercept payloads, extracts inputs, and selects corresponding prompt template based on features."),
        ("3. Prompt Synthesis", "Backend appends user parameters into highly engineered, detailed templates defining the expected output markdown structures."),
        ("4. Gemini API Inference", "Backend communicates with Google Gemini API via Client SDK libraries using user session API Key credentials."),
        ("5. Output Parsing & Render", "Frontend receives Gemini's markdown response, handles loading indicators toggles, and parses markdown into HTML page structures.")
    ]

    for step_num, (title, desc) in enumerate(flow_steps):
        # Draw a bullet for each step
        sp = tf_flow.add_paragraph()
        sp.text = f"{step_num + 1}. {title}"
        sp.font.name = 'Outfit'
        sp.font.size = Pt(16)
        sp.font.bold = True
        sp.font.color.rgb = COLOR_PURPLE
        sp.space_before = Pt(8)

        sd = tf_flow.add_paragraph()
        sd.text = desc
        sd.font.name = 'Inter'
        sd.font.size = Pt(14)
        sd.font.color.rgb = COLOR_WHITE
        sd.space_after = Pt(10)
        sd.level = 0
        sd.space_before = Pt(2)

    # ==========================================
    # SLIDE 7: Wireframes/Mock diagrams
    # ==========================================
    slide_7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_7)
    add_slide_title(slide_7, "Wireframes/Mock diagrams of the proposed solution")

    tb_mock = slide_7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_mock = tb_mock.text_frame
    tf_mock.word_wrap = True
    tf_mock.margin_left = tf_mock.margin_right = tf_mock.margin_top = tf_mock.margin_bottom = 0

    p = tf_mock.paragraphs[0]
    p.text = "Dashboard Interface Visual System Design"
    p.font.name = 'Outfit'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(16)

    mock_elements = [
        ("Unified Responsive Sidebar", "Left vertical navigation bar containing logo, active stats summaries, and page navigation routes links."),
        ("API Configuration Header", "Top bar with glowing key connector indicators (pulsing orange for setup, glowing green for active) connecting user sessions to Gemini."),
        ("Split-Screen Workspace Grid", "Two-column grid displaying input parameters form controls on the left, and glassmorphic recommendation panels on the right."),
        ("Actionable Output Panels", "Integrated clipboard copy, printer layouts exports, and document formatting rendering tools."),
        ("Real-time Shimmer loaders", "Animated loaders displaying feature-specific messages (e.g. 'Optimizing agenda...') during API processing.")
    ]

    for idx, (title, desc) in enumerate(mock_elements):
        mp = tf_mock.add_paragraph()
        mp.text = f"• {title}: "
        mp.font.name = 'Outfit'
        mp.font.size = Pt(16)
        mp.font.bold = True
        mp.font.color.rgb = COLOR_PURPLE
        mp.space_before = Pt(8)

        run = mp.add_run()
        run.text = desc
        run.font.name = 'Inter'
        run.font.size = Pt(15)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = False

    # ==========================================
    # SLIDE 8: Architecture diagram
    # ==========================================
    slide_8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_8)
    add_slide_title(slide_8, "Architecture diagram of the proposed solution")

    tb_arch = slide_8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_arch = tb_arch.text_frame
    tf_arch.word_wrap = True
    tf_arch.margin_left = tf_arch.margin_right = tf_arch.margin_top = tf_arch.margin_bottom = 0

    p = tf_arch.paragraphs[0]
    p.text = "LifeMate AI Architecture Blueprint"
    p.font.name = 'Outfit'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(20)

    arch_layers = [
        ("Layer 1: Client UI Interface", "HTML5 elements structure, CSS3 custom dark-theme glassmorphism system variables, and JavaScript AJAX router handlers."),
        ("Layer 2: Web Server Controller", "Python Flask server managing session-level API credentials variables, serving assets files, and controlling API endpoints."),
        ("Layer 3: prompt compiler", "Feature-specific prompt engineering templates designed to construct structured, context-rich queries for Gemini model processing."),
        ("Layer 4: Gemini AI Engine", "Google Gemini API (gemini-2.5-flash) executing reasoning tasks to synthesize schedules, budgets, plans, or emergency procedures."),
        ("Layer 5: Output Utilities", "Dynamic Markdown parser rendering, printable PDF formats generator engines, and browser clipboard utilities.")
    ]

    for idx, (layer, desc) in enumerate(arch_layers):
        ap = tf_arch.add_paragraph()
        ap.text = f"⚙ {layer}: "
        ap.font.name = 'Outfit'
        ap.font.size = Pt(16)
        ap.font.bold = True
        ap.font.color.rgb = COLOR_PURPLE
        ap.space_before = Pt(8)

        run = ap.add_run()
        run.text = desc
        run.font.name = 'Inter'
        run.font.size = Pt(15)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = False

    # ==========================================
    # SLIDE 9: Technologies Used
    # ==========================================
    slide_9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_9)
    add_slide_title(slide_9, "Technologies / Google / Nvidia Services used in the solution")

    tb_tech = slide_9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_tech = tb_tech.text_frame
    tf_tech.word_wrap = True
    tf_tech.margin_left = tf_tech.margin_right = tf_tech.margin_top = tf_tech.margin_bottom = 0

    p = tf_tech.paragraphs[0]
    p.text = "Why we chose our technical stack:"
    p.font.name = 'Outfit'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(20)

    techs = [
        ("Google Gemini API (gemini-2.5-flash)", "Core AI reasoning engine. We chose 2.5-flash for its low-latency performance, cost-efficiency, and advanced instruction-following capabilities across multiple topics."),
        ("Flask Framework (Python)", "Serves as the server backbone. Chosen for its lightweight footprint, fast routing setup, and compatibility with python-pptx and google-genai libraries."),
        ("google-genai Python SDK", "Leverages the latest, officially-supported modern Google GenAI Client libraries interfaces for prompt completions requests."),
        ("HTML5, CSS3 (Vanilla), JavaScript (ES6+)", "Ensures a fast, zero-dependency, and extremely responsive UI loading experience without heavy modern compiler frameworks."),
        ("python-pptx & python-dotenv", "Enables programmatic generation of hackathon presentations and secure loading of system-level environment credentials.")
    ]

    for idx, (tech_name, explanation) in enumerate(techs):
        tp = tf_tech.add_paragraph()
        tp.text = "✔ " + tech_name + ": "
        tp.font.name = 'Outfit'
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = COLOR_PURPLE
        tp.space_before = Pt(8)

        run = tp.add_run()
        run.text = explanation
        run.font.name = 'Inter'
        run.font.size = Pt(15)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = False

    # ==========================================
    # SLIDE 10: Snapshots of the prototype
    # ==========================================
    slide_10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_10)
    add_slide_title(slide_10, "Snapshots of the prototype")

    tb_snap = slide_10.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tf_snap = tb_snap.text_frame
    tf_snap.word_wrap = True
    tf_snap.margin_left = tf_snap.margin_right = tf_snap.margin_top = tf_snap.margin_bottom = 0

    p = tf_snap.paragraphs[0]
    p.text = "Visual Interface Snapshots Guide"
    p.font.name = 'Outfit'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(20)

    snap_guides = [
        ("Dashboard Console Panel", "Displays the stat widgets cards showing active tools status and quick shortcuts menu."),
        ("Daily Planner Generator Workspace", "Visualizes tasks inputs side-by-side with an optimized calendar output schedule formatted by Gemini AI."),
        ("Finances Diagnostics Screen", "Details monthly income vs expenditures checks and 50/30/20 diagnostic summaries suggestions."),
        ("Eco Blueprint Constructor Console", "Displays eco-habit suggestions checklists and neighbor composting hub blueprints."),
        ("Emergency Survival Protocol Screen", "Details evacuation steps and survival bags lists highlights in response to Flood/Fire selection.")
    ]

    for idx, (title, desc) in enumerate(snap_guides):
        sp = tf_snap.add_paragraph()
        sp.text = f"• {title}: "
        sp.font.name = 'Outfit'
        sp.font.size = Pt(16)
        sp.font.bold = True
        sp.font.color.rgb = COLOR_PURPLE
        sp.space_before = Pt(8)

        run = sp.add_run()
        run.text = desc + " (Run the local python server to capture fresh screenshots during the live demonstration)."
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = False

    # ==========================================
    # SLIDE 11: Thank You
    # ==========================================
    slide_11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide_11)

    tb_ty = slide_11.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(3.5))
    tf_ty = tb_ty.text_frame
    tf_ty.word_wrap = True
    tf_ty.margin_left = tf_ty.margin_right = tf_ty.margin_top = tf_ty.margin_bottom = 0

    p1 = tf_ty.paragraphs[0]
    p1.text = "Thank you!"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_CYAN
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf_ty.add_paragraph()
    p2.text = "LifeMate AI – Your digital companion for smarter, healthier, and more connected living."
    p2.font.name = 'Inter'
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_PURPLE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    p3 = tf_ty.add_paragraph()
    p3.text = "Powered by Google Cloud & Gemini AI"
    p3.font.name = 'Outfit'
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_MUTED
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(14)

    # Save presentation
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LifeMate_AI_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
